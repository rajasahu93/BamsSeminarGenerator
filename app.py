import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io

st.set_page_config(page_title="BAMS Seminar Presentation Generator", layout="centered")

st.markdown("""
    <h2 style='color: #2e7d32;'>🌿 BAMS Seminar Presentation Generator</h2>
    <p style='color: #555;'>Generating extensive, full-content academic slides loaded with deep core study text and complete textual framing.</p>
""", unsafe_allow_html=True)

topic = st.text_input("Enter Seminar Topic", "BENEFITS OF BREASTFEEDING ON OBSTETRICS")

def create_maximum_density_presentation(topic_name):
    prs = Presentation()
    
    # Custom slide dimensions for maximum page estate (Widescreen 16:9 layout format)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 15 detailed academic sections filled to maximum visual page length
    slides_data = [
        ("TITLE PRESENTATION", [
            f"COMPREHENSIVE ACADEMIC SEMINAR & CORE DISSERTATION",
            f"TOPIC: {topic_name}",
            "DEPARTMENT OF PRASUTI TANTRA EVAM STRIROGA & SAMHITA ADHYAYANA",
            "FACULTY OF AYURVEDIC MEDICINE AND SURGERY (BAMS)"
        ]),
        ("1. COMPREHENSIVE INTRODUCTION & BACKGROUND", [
            f"The subject of {topic_name} occupies a fundamental, core position in traditional Ayurvedic literature as well as contemporary neonatal and maternal healthcare frameworks.",
            "Ayurvedic science approaches this medical domain through holistic paradigms, emphasizing maternal nourishment, infant vitality, and the unhindered flow of life energy (Prana).",
            "Modern scientific validation is increasingly aligning with these ancient clinical observations, confirming profound immunological, biochemical, and physiological advantages across pediatric stages.",
            "The primary objective of this academic presentation is to bridge ancient scriptural references with rigorous, evidence-based modern correlations to establish a comprehensive clinical outlook."
        ]),
        ("2. ETYMOLOGY & CLASSICAL NIRUKTI", [
            "Linguistic derivation and etymology form the absolute cornerstone of Ayurvedic epistemology, uncovering the inherent functional attributes of complex physiological processes.",
            "Classical Sanskrit nomenclature is meticulously structured to express deep physiological mechanics and biochemical transformations rather than simple descriptive labels.",
            "Analyzing the root terms, suffixes, and contextual interpretations utilized throughout classical Samhitas sheds critical light on ancient conceptual frameworks.",
            "This etymological depth ensures that clinical applications remain entirely true to foundational classical principles while successfully communicating with modern medical terminology."
        ]),
        ("3. SAMHITA VACHANA & CLASSICAL REFERENCES", [
            "A comprehensive review of primary classical treatises including Charaka Samhita, Sushruta Samhita, and Ashtanga Hridaya reveals explicit scriptural citations.",
            "Specific passages located across Sutra Sthana, Sharira Sthana, and Uttara Tantra detail the precise qualities of pure breast milk (Stanya) and physiological production rules.",
            "Prominent classical commentators such as Chakrapani Datta and Dalhana provide profound insights regarding how maternal Rasa dhatu quality directly governs infant development.",
            "These scriptural citations establish an unbroken historical lineage of medical wisdom, validating the timeless efficacy and relevance of Ayurvedic obstetrics and pediatrics."
        ]),
        ("4. FUNDAMENTAL SIDDHANTA & BASIC PRINCIPLES", [
            "Ayurvedic physiology views lactation as an advanced, highly specialized tissue transformation driven by superior nourishment derived from wholesome maternal diet (Ahara).",
            "The delicate internal equilibrium of Tridosha—specifically Vyana Vata governing circulation, Sadhaka Pitta managing metabolism, and Tarpaka Kapha providing structural fluid—is vital.",
            "Any functional imbalance in Agni (digestive fire) or accumulation of Ama (endotoxins) at the tissue level compromises Stanya quality, triggering downstream pediatric issues.",
            "Mastery of these foundational Siddhantas empowers practitioners to execute root-cause-oriented interventions that resolve clinical pathologies effectively."
        ]),
        ("5. SAMPRAPTI & PATHOGENESIS MECHANICS", [
            "Pathogenesis in lactation and related obstetric conditions follows a well-defined sequence through classical stages, starting from improper dietary intake to tissue depletion.",
            "Vitiation of Prana Vata and Udana Vata via emotional stress or faulty lifestyle habits obstructs the smooth downward and outward flow pathways required for milk ejection.",
            "The pathological interaction between vitiated Doshas and local Dhatus causes structural and functional anomalies in Stanya production, presenting as clinical insufficiency.",
            "Early identification of these distinct pathological milestones allows practitioners to institute timely corrective herbal and dietary management protocols."
        ]),
        ("6. CLASSIFICATION & CLINICAL TYPOLOGY (BHEDA)", [
            "Classical texts classify lactation abnormalities and related obstetric parameters into rigorous typologies based on specific Dosha predominance and severity indices.",
            "Stanya Dushti (vitiations of breast milk) is systematically categorized into Vataja, Pittaja, Kaphaja, and Sannipatika variants, each displaying unique symptom profiles.",
            "Recognizing these specific classifications is critical for tailoring precise line-of-treatment strategies rather than relying on generalized therapeutic measures.",
            "Differential diagnostic parameters outlined in Samhitas help clinicians accurately distinguish between primary glandular insufficiencies and secondary lifestyle failures."
        ]),
        ("7. COMPREHENSIVE DIAGNOSTIC PROTOCOL (NIDANA PANCHAKA)", [
            "Clinical evaluation in Ayurveda utilizes the comprehensive five-fold framework of Nidana Panchaka to map the complete trajectory of maternal-infant health conditions.",
            "Nidana identifies root etiology such as improper maternal diet, emotional disturbance, and physical overexertion that disrupt normal postpartum physiological cycles.",
            "Purvaroopa identifies subtle prodromal signs like mammary heaviness or early milk supply drops before full clinical manifestations become apparent.",
            "Roopa, Upashaya, and Samprapti collectively complete the diagnostic matrix, empowering clinicians to establish a precise therapeutic diagnosis."
        ]),
        ("8. MANAGEMENT & CHIKITSA SIDDHANTA", [
            "Strategic management of conditions related to {topic_name} centers on restoring metabolic harmony, clearing bodily channels, and augmenting maternal vitality.",
            "The core Chikitsa Sutra emphasizes Deepana and Pachana measures to correct digestive fire prior to administering specialized galactagogue formulations (Stanyajanana).",
            "Internal administration of classical preparations enriched with nutritive herbs supports tissue regeneration and ensures optimal milk composition for the infant.",
            "Integrative therapeutic planning combines restorative dietary habits with psychological reassurance to maximize overall clinical effectiveness and recovery."
        ]),
        ("9. PATHYA-APATHYA & DIETETIC REGIMENS", [
            "Strict dietetic discipline (Pathya) is crucial during the postpartum and lactation period to sustain high-grade milk production without depleting maternal reserves.",
            "Recommended dietary articles include easily digestible, nutrient-dense foods cooked with traditional carminative spices that stimulate digestive fire and balance Vata.",
            "Prohibited habits and foods (Apathya) such as excessive dry, cold, or highly processed items are strictly avoided to prevent Ama accumulation and channel blockage.",
            "Adherence to structured daily routines (Dinacharya) ensures sustained maternal energy levels and robust immune defense throughout the nursing timeline."
        ]),
        ("10. MODERN SCIENTIFIC CORRELATION", [
            "Contemporary medical science strongly corroborates classical Ayurvedic principles regarding the profound biochemical and immunological superiority of natural feeding.",
            "Clinical research highlights live leukocytes, secretory antibodies, and bioactive growth factors in milk that mirror classical concepts of Ojas and immunity.",
            "Endocrinological studies validate the role of oxytocin and prolactin release during nursing, aligning directly with ancient neuro-hormonal and energetic pathways.",
            "Bridging these perspectives establishes a robust scientific foundation for integrating traditional wisdom into modern international obstetric guidelines."
        ]),
        ("11. PHARMACOLOGY & DRUG ACTION (DRAVYA GUNA)", [
            "Pharmacological evaluation of galactagogue herbs relies on understanding their unique Rasa, Guna, Virya, Vipaka, and Prabhava profiles.",
            "Substances possessing Madhura (sweet) and Snigdha (unctuous) properties act synergistically to nourish Rasa dhatu and stimulate mammary receptor sites.",
            "Cellular-level actions focus on enhancing metabolic assimilation and supporting systemic micro-circulation to facilitate unobstructed glandular secretion.",
            "Standardized quality control and bio-availability assessments ensure safety, purity, and predictable therapeutic outcomes in routine clinical practice."
        ]),
        ("12. CLINICAL OBSERVATION & CASE STUDY FRAMEWORK", [
            "Structured clinical observation frameworks provide empirical validation for classical Ayurvedic interventions applied within modern obstetric settings.",
            "Documenting patient cohorts undergoing traditional management reveals statistically significant improvements in milk volume and infant growth metrics.",
            "Standardized scoring systems assess parameters such as maternal fatigue levels, emotional well-being, and digestive comfort across postnatal recovery phases.",
            "These findings contribute valuable data to the broader academic community, encouraging further evidence-based exploration of traditional healthcare."
        ]),
        ("13. PREVENTIVE ASPECTS & PUBLIC HEALTH (SWASTHAVRITTA)", [
            "Preventive healthcare principles embedded in Swasthavritta emphasize proactive management of maternal health to avert long-term chronic complications.",
            "Promoting natural lactation protocols serves as a primary public health intervention against infant malnutrition, allergies, and metabolic disorders.",
            "Rejuvenation (Rasayana) therapies administered during the postnatal period secure long-term vitality, preventing premature aging and constitutional depletion.",
            "Community-level awareness campaigns rooted in classical traditions play an essential role in improving overall societal health metrics."
        ]),
        ("14. CONCLUSION & FINAL SUMMARY", [
            "The systematic academic exploration of {topic_name} demonstrates the remarkable depth, precision, and enduring relevance of classical Ayurveda.",
            "Synthesizing scriptural wisdom with modern scientific correlations confirms that traditional paradigms offer sophisticated solutions for contemporary health challenges.",
            "Future research directions must focus on large-scale clinical trials and standardized molecular evaluations to further validate ancient medical doctrines.",
            "The ultimate clinical objective remains the holistic well-being of both mother and child through balanced, compassionate, and evidence-based care."
        ]),
        ("15. REFERENCES & BIBLIOGRAPHY", [
            "1. Charaka Samhita, Ayurveda Dipika Commentary by Chakrapani Datta, Chaukhamba Surbharati Prakashan, Varanasi.",
            "2. Sushruta Samhita, Nibandha Samgraha Commentary by Dalhana, edited by Vaidya Jadavaji Trikamji Acharya, Chaukhamba Orientalia.",
            "3. Ashtanga Hridaya of Vagbhata, Vidyotini Hindi Commentary by Kaviraj Atrideva Gupta, Chaukhamba Sanskrit Series Office.",
            "4. Contemporary peer-reviewed journals in Obstetrics, Gynecology, and International Integrative Ayurvedic Medicine databases."
        ])
    ]

    for title, paragraphs in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout for maximum control
        
        # Add Header background banner area
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.1))
        tf_h = header_box.text_frame
        tf_h.word_wrap = True
        p_h = tf_h.paragraphs[0]
        p_h.text = title
        p_h.font.size = Pt(24)
        p_h.font.bold = True
        p_h.font.color.rgb = RGBColor(46, 125, 50) # Forest green BAMS theme
        
        # Add main content text frame covering the entire lower page area
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.733), Inches(5.3))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for i, para_text in enumerate(paragraphs):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = "• " + para_text
            p.font.size = Pt(14.5)  # Large, highly legible full-page font size
            p.font.color.rgb = RGBColor(30, 30, 30)
            p.space_after = Pt(12)  # Generous spacing to fill the page length cleanly
            
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

if st.button("Generate Full-Page 16:9 Detailed Seminar Presentation"):
    with st.spinner("Rendering full-page expanded layout across all slides..."):
        try:
            pptx_file = create_maximum_density_presentation(topic)
            st.success("Successfully generated full-page slides!")
            st.download_button(
                label="📥 Download Full-Page Presentation (.pptx)",
                data=pptx_file,
                file_name=f"{topic.replace(' ', '_')}_FullPage_Seminar.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        except Exception as e:
            st.error(f"Error generating presentation: {e}")
